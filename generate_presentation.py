import csv
import datetime
import glob
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Constants
OUTPUT_DIR = "pptx"
FIGURES_DIR = "output/figures"
DATA_DIR = "output"
DATE_STR = datetime.datetime.now().strftime("%Y%m%d")
VERSION = "v3"
OUTPUT_FILE = os.path.join(
    OUTPUT_DIR, f"Health_Economic_Analysis_{VERSION}_{DATE_STR}.pptx"
)


def create_presentation():  # noqa: C901
    prs = Presentation()

    # Set slide dimensions (16:9 aspect ratio is default, but let's ensure)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Health Economic Analysis: Policy Interventions"
    subtitle.text = f"New Zealand Medical Journal Submission\nGenerated: {datetime.datetime.now().strftime('%Y-%m-%d')}\nVersion: {VERSION}"

    # 2. Outline Slide
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Outline"
    content = slide.placeholders[1]
    content.text = (
        "1. Decision Model Structures\n"
        "2. Cost-Effectiveness Analysis Results\n"
        "3. Comparative ICER Summary\n"
        "4. Probabilistic Sensitivity Analysis\n"
        "5. Distributional Cost-Effectiveness Analysis (DCEA)\n"
        "6. Budget Impact Analysis\n"
        "7. Value of Information Analysis"
    )

    # Helper to add image slide
    def add_image_slide(img_path, slide_title, note_text=""):
        if not os.path.exists(img_path):
            print(f"Warning: Image not found: {img_path}")
            return

        # Use Title Only layout
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_shape = slide.shapes.title
        title_shape.text = slide_title

        # Add image
        # Calculate position to center and fit
        # Margins
        left = Inches(0.5)
        top = Inches(1.5)
        width = prs.slide_width - Inches(1.0)
        height = prs.slide_height - Inches(2.0)

        try:
            from PIL import Image

            with Image.open(img_path) as img:
                img_w, img_h = img.size
                aspect_ratio = img_w / img_h
        except ImportError:
            aspect_ratio = 1.33  # Fallback

        slide_ratio = width / height

        if aspect_ratio > slide_ratio:
            # Image is wider than slide area, fit to width
            pic = slide.shapes.add_picture(img_path, left, top, width=width)
            # Center vertically
            pic_h = width / aspect_ratio
            pic.top = int(top + (height - pic_h) / 2)
        else:
            # Image is taller, fit to height
            pic = slide.shapes.add_picture(img_path, left, top, height=height)
            # Center horizontally
            pic_w = height * aspect_ratio
            pic.left = int(left + (width - pic_w) / 2)

        # Add notes
        if note_text:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = note_text

    # Helper to add table slide
    def add_table_slide(csv_path, slide_title):
        if not os.path.exists(csv_path):
            print(f"Warning: CSV not found: {csv_path}")
            return

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        title_shape = slide.shapes.title
        title_shape.text = slide_title

        rows = []
        with open(csv_path) as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            return

        # Filter columns to fit slide (Intervention, Perspective, ICER, Cost-Effective)
        # Original headers: Intervention,Perspective,Cost New Treatment,Cost Standard Care,QALYs New Treatment,QALYs Standard Care,Incremental Cost,Incremental QALYs,ICER ($/QALY),Net Monetary Benefit,Cost-Effective (WTP=$50k)
        # Indices: 0, 1, 8, 10 (approx)

        header = rows[0]
        target_cols = [
            "Intervention",
            "Perspective",
            "ICER ($/QALY)",
            "Cost-Effective (WTP=$50k)",
        ]
        col_indices = [i for i, h in enumerate(header) if h in target_cols]

        # Table dimensions
        rows_count = len(rows)
        cols_count = len(col_indices)

        left = Inches(0.5)
        top = Inches(1.5)
        width = prs.slide_width - Inches(1.0)
        height = prs.slide_height - Inches(2.0)

        table = slide.shapes.add_table(
            rows_count, cols_count, left, top, width, height
        ).table

        # Set column widths
        # table.columns[0].width = Inches(4.0)

        for r_idx, row in enumerate(rows):
            for c_idx, col_idx in enumerate(col_indices):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(row[col_idx])

                # Format header
                if r_idx == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark Blue
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    paragraph.font.bold = True
                else:
                    # Format numbers
                    try:
                        val = float(row[col_idx])
                        if "ICER" in header[col_idx]:
                            cell.text = f"${val:,.0f}"
                    except ValueError:
                        pass

                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(10)  # Smaller font for data

    # 3. Add Slides

    # Decision Trees
    trees = sorted(glob.glob(os.path.join(FIGURES_DIR, "decision_tree_*.png")))
    for img in trees:
        name = (
            os.path.basename(img)
            .replace("decision_tree_", "")
            .replace(".png", "")
            .replace("_", " ")
            .title()
        )
        add_image_slide(
            img,
            f"Model Structure: {name}",
            f"Decision tree structure for {name} intervention.",
        )

    # CE Planes
    planes = sorted(glob.glob(os.path.join(FIGURES_DIR, "ce_plane_*.png")))
    for img in planes:
        name = (
            os.path.basename(img)
            .replace("ce_plane_", "")
            .replace(".png", "")
            .replace("_", " ")
            .title()
        )
        add_image_slide(
            img,
            f"Cost-Effectiveness Plane: {name}",
            "Scatter plot of incremental costs vs incremental QALYs.",
        )

    # Comparative ICER Table
    add_table_slide(
        os.path.join(DATA_DIR, "comparative_icer_table.csv"), "Comparative ICER Summary"
    )

    # Probabilistic Sensitivity Analysis (CEAC/CEAF)
    ceac = os.path.join(FIGURES_DIR, "probabilistic_analysis_ceac_ceaf.png")
    add_image_slide(
        ceac,
        "Probabilistic Sensitivity Analysis",
        "Cost-Effectiveness Acceptability Curve and Frontier.",
    )

    # Equity Analysis
    equity = os.path.join(FIGURES_DIR, "equity_analysis.png")
    add_image_slide(
        equity,
        "Equity Analysis",
        "Distributional impact of interventions across socioeconomic groups.",
    )

    dcea_wtp = os.path.join(FIGURES_DIR, "dcea_willingness_to_pay.png")
    add_image_slide(
        dcea_wtp,
        "DCEA: Willingness to Pay",
        "Willingness to pay for equity-weighted health gains.",
    )

    # Budget Impact Analysis
    bia = os.path.join(FIGURES_DIR, "BIA_budget_impact_analysis_dashboard.png")
    # Fallback if BIA prefix isn't there
    if not os.path.exists(bia):
        bia = os.path.join(FIGURES_DIR, "budget_impact_analysis_dashboard.png")
    add_image_slide(
        bia,
        "Budget Impact Analysis",
        "Financial implications of the interventions over time.",
    )

    # VOI
    evpi = os.path.join(FIGURES_DIR, "voi_analysis_evpi.png")
    add_image_slide(
        evpi,
        "Value of Information: EVPI",
        "Expected Value of Perfect Information analysis.",
    )

    # 4. End Slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Questions & Discussion"

    # Save
    if not os.path.exists(OUTPUT_DIR):
        os.makedirs(OUTPUT_DIR)
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")


if __name__ == "__main__":
    create_presentation()
